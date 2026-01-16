from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r"C:\Users\Mohid\OneDrive\Documents\GitHub\Concept-Mapping-Concept-Extraction\PythonPPTX_Testing\test.pptx"
txt_path = r"C:\Users\Mohid\OneDrive\Documents\GitHub\Concept-Mapping-Concept-Extraction\PythonPPTX_Testing\extractionTest.txt"

prs = Presentation(pptx_path)

def get_alt_text(shape):
    try:
        cNvPr = shape._element.xpath(".//*[local-name() = 'cNvPr']")
        if cNvPr:
            title = cNvPr[0].get("title")
            descr = cNvPr[0].get("descr")
            if descr and descr.strip():
                return descr.strip()
            if title and title.strip():
                return title.strip()
    except Exception:
        pass
    return ""

with open(txt_path, "w", encoding = "utf-8") as file:
    for slide_num, slide in enumerate(prs.slides, start = 1):
        file.write(f"--- Slide {slide_num} ---\n")

        for shape in slide.shapes:
            # Text boxes
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        file.write(text + "\n")

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                alt_text = get_alt_text(shape)
                if alt_text:
                    file.write(f"[Image Alt Text] {alt_text}\n")
                else:
                    file.write("[Image Missing Alt Text]\n")

        file.write("\n")

print("Extraction complete.")
