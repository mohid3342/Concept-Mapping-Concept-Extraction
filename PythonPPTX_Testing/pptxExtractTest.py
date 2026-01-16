from pptx import Presentation

pptx_path = r"C:\Users\Mohid\OneDrive\Documents\GitHub\MohidCode\MainCode\Code\Python\VSCode\COSC 481\test.pptx"
txt_path = r"C:\Users\Mohid\OneDrive\Documents\GitHub\MohidCode\MainCode\Code\Python\VSCode\COSC 481\extracted_text.txt"

prs = Presentation(pptx_path)

with open(txt_path, "w", encoding="utf-8") as file:
    for slide_num, slide in enumerate(prs.slides, start=1):
        file.write(f"--- Slide {slide_num} ---\n")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    file.write(text + "\n")
        file.write("\n")

print("Text extraction complete. File saved as extracted_text.txt")
