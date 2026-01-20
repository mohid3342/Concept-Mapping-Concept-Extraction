from pptx import Presentation
import json
from pptx.enum.shapes import MSO_SHAPE_TYPE
import base64
from openai import OpenAI


client = OpenAI(
    api_key="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"
)

def image_to_base64(image):
    return base64.b64encode(image.blob).decode("utf-8")


def describe_image_with_ai(image_base64):
    response = client.responses.create(
        model="gpt-4o",
        input=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "input_text",
                        "text": "Describe this image for a college lecture slide in 25 words"
                    },
                    {
                        "type": "input_image",
                        "image_url": f"data:image/jpeg;base64,{image_base64}",
                        "detail": "low"
                    }
                ]
            }
        ]
    )
    return response.output_text


def get_alt_text(shape):
    try:
        cNvPr = shape._element.xpath(".//*[local-name() = 'cNvPr']")
        if cNvPr:
            title = cNvPr[0].get("title")
            descr = cNvPr[0].get("descr")
            if descr and descr.strip():
                return descr.strip()
            if title and title.strip():
                return title.strip()
    except Exception:
        pass
    return ""

def extract_pptx_data(file_path):
    prs = Presentation(file_path)
    ppt_data = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "slide_number": slide_index,
            "title": slide.shapes.title.text if slide.shapes.title else None,
            "text_blocks": [],
            "images": []
        }
        for shape in slide.shapes:
            # Text shapes
            if shape.has_text_frame:
                text_block = []
                for paragraph in shape.text_frame.paragraphs:
                    para_content = {
                        "text": paragraph.text,
                        "runs": []
                    }
                    for run in paragraph.runs:
                        run_info = {
                            "text": run.text,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "underline": run.font.underline,
                            "font_name": run.font.name,
                            "font_size": run.font.size.pt if run.font.size else None
                        }
                        para_content["runs"].append(run_info)
                    text_block.append(para_content)
                slide_info["text_blocks"].append(text_block)

            # Images with alt text
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # Picture
                try:
                    alt_text = get_alt_text(shape)
                    image_base64 = image_to_base64(shape.image)
                    ai_description = describe_image_with_ai(image_base64)

                    slide_info["images"].append({
                        "alt_text": alt_text,
                        "ai_description": ai_description
                    })
                except Exception as e:
                    slide_info["images"].append({
                        "alt_text": alt_text if alt_text else "No alt text",
                        "ai_description": None,
                        "error": str(e)
                    })
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            slide_info["notes"] = notes_text
        ppt_data.append(slide_info)
        
 

    return ppt_data

# -----------------------------
# Example usage
file_path = "test.pptx"
ppt_structure = extract_pptx_data(file_path)

# Save to a JSON file
output_file = "ppt_test_data.json"
with open(output_file, "w", encoding="utf-8") as f:
    json.dump(ppt_structure, f, indent=2, ensure_ascii=False)

print(f"Data saved to {output_file}")

